import PyPDF2
import docx2txt
import pptx


def extract_text_from_word(file_path):
    text = docx2txt.process(file_path)
    return text


def extract_text_from_pdf(file_path):
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        text = ''
        for page in reader.pages:
            text += page.extract_text()
    return text


def extract_text_from_powerpoint(file_path):
    prs = pptx.Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text
