import os
import docx2txt
import PyPDF2
import pptx
from database import create_table
from filter_stopwords import filter_stopwords

def extract_text_from_word(file_path):
    text = docx2txt.process(file_path)
    return text

def extract_text_from_pdf(file_path):
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        text = ''
        for page in reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_powerpoint(file_path):
    prs = pptx.Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def build_reverse_index(repo_path):
    create_table()
    reverse_index = {}

    for root, dirs, files in os.walk(repo_path):
        for file_name in files:
            file_path = os.path.join(root, file_name)
            file_extension = os.path.splitext(file_name)[1].lower()
            if file_extension == '.docx':
                try:
                    text = extract_text_from_word(file_path)
                    filtered_text = filter_stopwords(text)
                    for word in filtered_text.split():
                        if word in reverse_index:
                            reverse_index[word].append(file_path)
                        else:
                            reverse_index[word] = [file_path]
                except Exception as e:
                    print(f"Skipped non-Word file: {file_path} ({str(e)})")
            elif file_extension == '.pdf':
                try:
                    text = extract_text_from_pdf(file_path)
                    filtered_text = filter_stopwords(text)
                    for word in filtered_text.split():
                        if word in reverse_index:
                            reverse_index[word].append(file_path)
                        else:
                            reverse_index[word] = [file_path]
                except (PyPDF2.utils.PdfReadError, ValueError) as e:
                    print(f"Skipped non-PDF file: {file_path} ({str(e)})")
            elif file_extension == '.pptx':
                try:
                    text = extract_text_from_powerpoint(file_path)
                    filtered_text = filter_stopwords(text)
                    for word in filtered_text.split():
                        if word in reverse_index:
                            reverse_index[word].append(file_path)
                        else:
                            reverse_index[word] = [file_path]
                except Exception as e:
                    print(f"Skipped non-PowerPoint file: {file_path} ({str(e)})")

    conn = sqlite3.connect('reverse_index.db')
    cursor = conn.cursor()

    for word, files in reverse_index.items():
        files_str = ' '.join(files)
        cursor.execute('INSERT INTO reverse_index VALUES (?, ?)', (word, files_str))

    conn.commit()
    conn.close()
